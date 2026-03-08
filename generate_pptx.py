import subprocess
import sys
import os

try:
    import pptx
except ImportError:
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
    except Exception as e:
        print("Could not install python-pptx:", e)
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ALL SAFE SOC (v3.0.0)"
subtitle.text = "Advanced Cyber Threat Intelligence & Security Operations Center Platform\nPrototype Overview"

# Slide 2: Core Overview
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "1. Core Overview"
tf = body_shape.text_frame
tf.text = "ALL SAFE SOC is a comprehensive, multi-source cybersecurity platform."
p = tf.add_paragraph()
p.text = "Aggregates, analyzes, and visualizes real-time cyber threats."
p.level = 1
p = tf.add_paragraph()
p.text = "Modern backend built with Python and FastAPI."
p.level = 1
p = tf.add_paragraph()
p.text = "Integrates top-tier threat intelligence APIs, OSINT tools, and AI models."
p.level = 1

# Slide 3: Real-Time Global Threat Detection
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "2. Key Features: Threat Detection"
tf = body_shape.text_frame
tf.text = "Real-Time Global Threat Detection (Live Attack Map)"
p = tf.add_paragraph()
p.text = "Aggregates live threat feeds from multiple global sources."
p.level = 1
p = tf.add_paragraph()
p.text = "Streams data instantly to the frontend via WebSockets."
p.level = 1
p = tf.add_paragraph()
p.text = "Simulated Honeybot Network for dynamic visualization of attacks."
p.level = 1

# Slide 4: Multi-Source Threat Engine
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "3. Multi-Source Threat Engine"
tf = body_shape.text_frame
tf.text = "Beast Mode Integrations"
p = tf.add_paragraph()
p.text = "AlienVault OTX: Global indicator feeds and malicious pulses."
p.level = 1
p = tf.add_paragraph()
p.text = "AbuseIPDB: Community-driven IP abuse confidence scoring."
p.level = 1
p = tf.add_paragraph()
p.text = "ThreatFox: IOCs for malware and botnets."
p.level = 1
p = tf.add_paragraph()
p.text = "HoneyDB: Real-time honeypot hits."
p.level = 1
p = tf.add_paragraph()
p.text = "VirusTotal: Core Engine for deep static/dynamic analysis."
p.level = 1

# Slide 5: AI-Powered Analysis
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "4. AI-Powered Analysis"
tf = body_shape.text_frame
tf.text = "Intelligent Threat Summarization"
p = tf.add_paragraph()
p.text = "Primary AI: Google Gemini (2.5 Flash / 2.0 Flash) with key rotation."
p.level = 1
p = tf.add_paragraph()
p.text = "Fallback AI: Groq (LLaMA 3.3 70B / 3.1 8B)."
p.level = 1
p = tf.add_paragraph()
p.text = "Automatically generates professional, concise threat summaries and risk assessments."
p.level = 1

# Slide 6: Comprehensive Scanning Suite
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "5. Comprehensive Scanning Suite"
tf = body_shape.text_frame
tf.text = "Scanner Capabilities"
p = tf.add_paragraph()
p.text = "IP Scanner: OTX, AbuseIPDB, ThreatFox, VirusTotal, Shodan, WHOIS, Nmap."
p.level = 1
p = tf.add_paragraph()
p.text = "URL, Domain, File & Hash Scanners."
p.level = 1
p = tf.add_paragraph()
p.text = "Phone OSINT Scanner (Spam scoring & global region parsing)."
p.level = 1
p = tf.add_paragraph()
p.text = "Identity / Breach Scanner (Data leak correlation)."
p.level = 1
p = tf.add_paragraph()
p.text = "Job Scam Detector (Regex patterns + AI reasoning)."
p.level = 1

# Slide 7: Tech Stack & Architecture
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "6. Tech Stack & Architecture"
tf = body_shape.text_frame
tf.text = "Modern & Scalable Architecture"
p = tf.add_paragraph()
p.text = "Backend Layer: Python 3, FastAPI, Uvicorn."
p.level = 1
p = tf.add_paragraph()
p.text = "Database: SQLite3 (Local scan history and telemetry logging)."
p.level = 1
p = tf.add_paragraph()
p.text = "Networking: Asyncio, HTTPX, WebSockets."
p.level = 1
p = tf.add_paragraph()
p.text = "AI Framework: google-genai SDK, Groq SDK."
p.level = 1

# Slide 8: Use Cases & Roadmap
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "7. Use Cases & Roadmap"
tf = body_shape.text_frame
tf.text = "Who is this for?"
p = tf.add_paragraph()
p.text = "SOC Analysts for rapid triage."
p.level = 1
p = tf.add_paragraph()
p.text = "Security teams for global situational awareness."
p.level = 1
p = tf.add_paragraph()
p.text = "Regular users verifying threats (job scams, phishing links)."
p.level = 1
p = tf.add_paragraph()
p.text = "Future Roadmap: Wazuh SIEM integration, automated incident response."
p.level = 1

out_path = os.path.join("c:/Users/mahen/OneDrive/Apps/NEXATHAN", "ALL_SAFE_Prototype.pptx")
prs.save(out_path)
print(f"Successfully generated {out_path}")
